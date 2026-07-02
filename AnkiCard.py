import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import os

class PixelDisplay(tk.Canvas):
    """像素显示屏组件"""
    def __init__(self, parent, width=60, height=18, pixel_size=3, bg_color="#2A2E33", pixel_color="#2A2E33", **kwargs):
        super().__init__(parent, width=width*pixel_size, height=height*pixel_size, bg=bg_color, bd=0, highlightthickness=0, **kwargs)
        self.pixel_size = pixel_size
        self.width = width
        self.height = height
        self.bg_color = bg_color
        self.pixel_color = pixel_color
        self.grid = [[False] * width for _ in range(height)]
        self.draw_grid()
    
    def draw_grid(self):
        """绘制网格背景"""
        self.delete("all")
        for y in range(self.height):
            for x in range(self.width):
                x1 = x * self.pixel_size
                y1 = y * self.pixel_size
                x2 = x1 + self.pixel_size
                y2 = y1 + self.pixel_size
                self.create_rectangle(x1, y1, x2, y2, fill=self.bg_color, outline="#3a3a3a", width=1)
    
    def set_pixel(self, x, y, state=True):
        """设置单个像素"""
        if 0 <= x < self.width and 0 <= y < self.height:
            self.grid[y][x] = state
            self.draw_pixels()
    
    def draw_pixels(self):
        """绘制所有像素"""
        self.delete("all")
        for y in range(self.height):
            for x in range(self.width):
                x1 = x * self.pixel_size
                y1 = y * self.pixel_size
                x2 = x1 + self.pixel_size
                y2 = y1 + self.pixel_size
                
                if self.grid[y][x]:
                    self.create_rectangle(x1, y1, x2, y2, fill=self.pixel_color, outline=self.pixel_color, width=0)
                else:
                    self.create_rectangle(x1, y1, x2, y2, fill=self.bg_color, outline="#3a3a3a", width=1)
    
    def clear(self):
        """清空显示"""
        self.grid = [[False] * self.width for _ in range(self.height)]
        self.draw_pixels()
    
    def display_text(self, text, row=0):
        """显示文本（简单像素字体）"""
        self.clear()
        col = 0
        for char in text:
            if col >= self.width - 4:
                break
            self.draw_char(char, col, row)
            col += 5
    
    def draw_char(self, char, start_col, row):
        """绘制单个字符的像素"""
        fonts = {
            '0': [[1,1,1], [1,0,1], [1,0,1], [1,0,1], [1,1,1]],
            '1': [[0,1,0], [1,1,0], [0,1,0], [0,1,0], [1,1,1]],
            '2': [[1,1,1], [0,0,1], [1,1,1], [1,0,0], [1,1,1]],
            '3': [[1,1,1], [0,0,1], [1,1,1], [0,0,1], [1,1,1]],
            '4': [[1,0,1], [1,0,1], [1,1,1], [0,0,1], [0,0,1]],
            '5': [[1,1,1], [1,0,0], [1,1,1], [0,0,1], [1,1,1]],
            '6': [[1,1,1], [1,0,0], [1,1,1], [1,0,1], [1,1,1]],
            '7': [[1,1,1], [0,0,1], [0,1,0], [0,1,0], [0,1,0]],
            '8': [[1,1,1], [1,0,1], [1,1,1], [1,0,1], [1,1,1]],
            '9': [[1,1,1], [1,0,1], [1,1,1], [0,0,1], [1,1,1]],
            ' ': [[0,0,0], [0,0,0], [0,0,0], [0,0,0], [0,0,0]],
        }
        
        pattern = fonts.get(char, fonts[' '])
        for y, line in enumerate(pattern):
            for x, pixel in enumerate(line):
                if pixel and row + y < self.height:
                    self.set_pixel(start_col + x, row + y, True)

class EnglishFlashcardApp:
    MAX_CARDS = 100

    def __init__(self, root):
        self.root = root
        self.root.title("Anki卡片生成 - LED显示屏版")
        self.root.geometry("600x600")
        self.root.config(bg="#2A2E33")
        self.cards = []
        
        self.led_bg = "#2A2E33"
        self.led_color = "#FAEBF1"
        self.led_color_accent = "#FAEBF1"
        self.led_color_green = "#22c55e"
        self.led_color_red = "#ef4444"
        self.led_color_blue = "#0088aa"
        
        self.setup_ui()
    
    def setup_ui(self):
        menubar = tk.Menu(self.root)
        self.root.config(menu=menubar)
        
        file_menu = tk.Menu(menubar, tearoff=0, font=("Courier New", 12))
        menubar.add_cascade(label="FILE", menu=file_menu)
        file_menu.add_command(label="LOAD (TXT/MD)", command=self.load_from_file)
        file_menu.add_separator()
        file_menu.add_command(label="EXIT", command=self.root.quit)
        
        title_frame = tk.Frame(self.root, bg=self.led_bg)
        title_frame.pack(fill=tk.X, padx=9, pady=9)
        
        tk.Label(
            title_frame,
            text="█ ANKI CARD GENERATOR █",
            font=("Courier New", 17, "bold"),
            bg=self.led_bg,
            fg=self.led_color_accent
        ).pack()
        
        count_frame = tk.Frame(self.root, bg=self.led_bg)
        count_frame.pack(fill=tk.X, padx=10, pady=5)
        
        tk.Label(count_frame, text="CARDS:", font=("Courier New", 13, "bold"), bg=self.led_bg, fg=self.led_color).pack(side=tk.LEFT, padx=5)
        self.card_count_display = PixelDisplay(count_frame, width=10, height=6, pixel_size=6, bg_color="#2A2E33", pixel_color=self.led_color_accent)
        self.card_count_display.pack(side=tk.LEFT, padx=5)
        
        tk.Frame(self.root, height=2, bg=self.led_color).pack(fill=tk.X, padx=10, pady=8)
        
        input_frame = tk.LabelFrame(
            self.root,
            text="INPUT",
            font=("Courier New", 13, "bold"),
            bg=self.led_bg,
            fg=self.led_color,
            bd=2,
            relief=tk.RIDGE
        )
        input_frame.pack(fill=tk.X, padx=10, pady=8)
        
        tk.Label(input_frame, text="ENGLISH:", font=("Courier New", 12, "bold"), bg=self.led_bg, fg=self.led_color).pack(anchor="w", padx=12, pady=(10, 5))
        
        self.english_entry = tk.Entry(
            input_frame,
            font=("Courier New", 11),
            width=70,
            relief=tk.FLAT,
            bd=0,
            bg="#1a1e22",
            fg=self.led_color_accent,
            insertbackground=self.led_color_accent
        )
        self.english_entry.pack(padx=12, pady=5, ipady=6)
        
        tk.Label(input_frame, text="CHINESE:", font=("Courier New", 12, "bold"), bg=self.led_bg, fg=self.led_color).pack(anchor="w", padx=12, pady=(10, 5))
        
        self.chinese_entry = tk.Entry(
            input_frame,
            font=("Courier New", 11),
            width=70,
            relief=tk.FLAT,
            bd=0,
            bg="#1a1e22",
            fg=self.led_color_accent,
            insertbackground=self.led_color_accent
        )
        self.chinese_entry.pack(padx=12, pady=(0, 12), ipady=6)
        
        inner_button_frame = tk.Frame(self.root, bg=self.led_bg)
        inner_button_frame.pack(pady=10)
        
        save_btn = self.create_led_button(inner_button_frame, "SAVE", self.save_card, self.led_color_green)
        save_btn.pack(side=tk.LEFT, padx=6)
        
        generate_btn = self.create_led_button(inner_button_frame, "GEN PPT", self.generate_ppt, self.led_color_red)
        generate_btn.pack(side=tk.LEFT, padx=6)
        
        clear_btn = self.create_led_button(inner_button_frame, "CLEAR", self.clear_inputs, self.led_color_blue)
        clear_btn.pack(side=tk.LEFT, padx=6)
        
        tk.Label(
            self.root,
            text="▓ SAVED CARDS ▓",
            font=("Courier New", 13, "bold"),
            bg=self.led_bg,
            fg=self.led_color
        ).pack(anchor="w", padx=18, pady=(8, 5))
        
        display_frame = tk.Frame(self.root, bg=self.led_bg)
        display_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=(0, 10))
        
        self.card_display = scrolledtext.ScrolledText(
            display_frame,
            height=6,
            width=80,
            font=("Courier New", 11),
            relief=tk.FLAT,
            bd=0,
            bg="#1a1e22",
            fg=self.led_color,
            insertbackground=self.led_color
        )
        self.card_display.pack(fill=tk.BOTH, expand=True, padx=0, pady=0)
        self.card_display.config(state=tk.DISABLED)
    
    def create_led_button(self, parent, text, command, color):
        btn = tk.Button(
            parent,
            text=text,
            command=command,
            font=("Courier New", 15, "bold"),
            bg=self.led_bg,
            fg=color,
            padx=16,
            pady=6,
            relief=tk.RAISED,
            bd=1,
            activebackground="#2A2E33",
            activeforeground=self.led_color_accent,
            cursor="hand2",
            highlightthickness=1,
            highlightcolor=color
        )
        
        def on_enter(event):
            btn.config(relief=tk.SUNKEN, bd=2, fg=self.led_color_accent)
        
        def on_leave(event):
            btn.config(relief=tk.RAISED, bd=1, fg=color)
        
        btn.bind("<Enter>", on_enter)
        btn.bind("<Leave>", on_leave)
        
        return btn
    
    def save_card(self):
        english = self.english_entry.get().strip()
        chinese = self.chinese_entry.get().strip()
        
        if not english or not chinese:
            messagebox.showwarning("WARNING", "PLEASE INPUT ENGLISH AND CHINESE")
            return
        
        if len(self.cards) >= self.MAX_CARDS:
            messagebox.showwarning("WARNING", f"MAX {self.MAX_CARDS} CARDS ONLY")
            return
        
        self.cards.append({"english": english, "chinese": chinese})
        self.english_entry.delete(0, tk.END)
        self.chinese_entry.delete(0, tk.END)
        self.english_entry.focus()
        
        self.update_display()
    
    def update_display(self):
        count_str = str(len(self.cards)).zfill(2)
        self.card_count_display.display_text(count_str, row=0)
        
        self.card_display.config(state=tk.NORMAL)
        self.card_display.delete(1.0, tk.END)
        
        for i, card in enumerate(self.cards, 1):
            self.card_display.insert(tk.END, f"[CARD {i:02d}]\n", "header")
            self.card_display.insert(tk.END, f"ENG: {card['english']}\n", "content")
            self.card_display.insert(tk.END, f"CHN: {card['chinese']}\n", "content")
            self.card_display.insert(tk.END, "─" * 70 + "\n", "content")
        
        self.card_display.tag_config("header", font=("Courier New", 11))
        self.card_display.tag_config("content", font=("Courier New", 11))
        
        self.card_display.config(state=tk.DISABLED)
    
    def load_from_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("Text files", "*.txt"), ("Markdown files", "*.md")])
        if not file_path:
            return
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            pattern = r"===\s*第\d+页\s*===\s*英文:\s*(.+?)\s*中文:\s*(.+?)(?===|$)"
            matches = re.findall(pattern, content, re.DOTALL)
            
            if not matches:
                messagebox.showerror("ERROR", "FILE FORMAT ERROR: NO CARD DATA FOUND")
                return
            
            if len(matches) > self.MAX_CARDS:
                messagebox.showwarning("WARNING", f"INPUT EXCEEDS {self.MAX_CARDS} CARDS, ONLY FIRST {self.MAX_CARDS} WILL BE USED")
                matches = matches[:self.MAX_CARDS]
            
            self.cards = []
            for english, chinese in matches:
                self.cards.append({
                    "english": english.strip(),
                    "chinese": chinese.strip()
                })
            
            self.update_display()
            messagebox.showinfo("SUCCESS", f"LOADED {len(self.cards)} CARDS")
        
        except Exception as e:
            messagebox.showerror("ERROR", f"LOAD FAILED: {str(e)}")
    
    def generate_ppt(self):
        if not self.cards:
            messagebox.showwarning("WARNING", "NO CARDS TO GENERATE")
            return
        
        now = datetime.now()
        filename = now.strftime("Anki%y%m%d%H%M.pptx")
        file_path = os.path.join(os.getcwd(), filename)
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for idx, card in enumerate(self.cards, 1):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                left = Inches(0.8)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(2)
                
                text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = card['english']
                p.font.name = "Courier New"
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 0.9
                
                left = Inches(0.8)
                top = Inches(5.2)
                width = Inches(9)
                height = Inches(2)
                
                text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = card['chinese']
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(38)
                p.font.color.rgb = RGBColor(25, 25, 112)
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 0.9
            
            prs.save(file_path)
            
            txt_path = file_path.replace('.pptx', '.txt')
            with open(txt_path, 'w', encoding='utf-8') as f:
                for idx, card in enumerate(self.cards, 1):
                    f.write(f"=== 第{idx}页 ===\n")
                    f.write(f"英文: {card['english']}\n")
                    f.write(f"中文: {card['chinese']}\n\n")
            
            messagebox.showinfo("SUCCESS", f"PPT GENERATED\nPATH: {file_path}")
        
        except Exception as e:
            messagebox.showerror("ERROR", f"GENERATION FAILED: {str(e)}")
    
    def clear_inputs(self):
        self.english_entry.delete(0, tk.END)
        self.chinese_entry.delete(0, tk.END)
        self.english_entry.focus()

if __name__ == "__main__":
    root = tk.Tk()
    app = EnglishFlashcardApp(root)
    root.mainloop()